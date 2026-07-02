import os
import csv
import io
import uuid
from datetime import datetime, date, timedelta
from flask import Flask, render_template, request, redirect, url_for, flash, send_file, Response
from flask_login import LoginManager, login_user, logout_user, login_required, current_user
from werkzeug.utils import secure_filename
from sqlalchemy import extract, and_

# ReportLab imports for PDF generation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

from models import db, User, Transaction, Budget, UpcomingBill, CustomCategory

app = Flask(__name__)
app.secret_key = 'expense-tracker-secret-key-12345'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///database.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['MAX_CONTENT_LENGTH'] = 5 * 1024 * 1024  # 5MB size limit

db.init_app(app)

login_manager = LoginManager()
login_manager.login_view = 'login'
login_manager.init_app(app)

# Configure Uploads
UPLOAD_FOLDER = os.path.join(app.root_path, 'static', 'receipts')
PROFILE_PICS_FOLDER = os.path.join(app.root_path, 'static', 'profile_pics')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROFILE_PICS_FOLDER, exist_ok=True)
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'webp'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# --- Navigation Active Page Context Processor ---
@app.context_processor
def inject_active_page():
    return dict(active_page=request.endpoint)

TRANSLATIONS = {
    'en': {
        'dashboard': 'Dashboard',
        'transactions': 'Transactions',
        'reports': 'Reports',
        'settings': 'Settings',
        'total_balance': 'Total Balance',
        'monthly_income': 'Monthly Income',
        'monthly_expenses': 'Monthly Expenses',
        'monthly_savings': 'Monthly Savings',
        'savings_rate': 'Savings Rate',
        'budget_consumptions': 'Budget Consumptions',
        'highest_spending_category': 'Highest Spending Category',
        'spent_breakdown': 'Spent Breakdown',
        'spending_trends': 'Spending Trends',
        'upcoming_bills': 'Upcoming Bills',
        'recent_transactions': 'Recent Transactions',
        'add_transaction': 'Add Transaction',
        'set_budget': 'Set Budget',
        'add_bill': 'Add Bill',
        'export_options': 'Export Options',
        'select_month': 'Select Month',
        'export_format': 'Export Format',
        'pdf_document': 'PDF Document',
        'csv_spreadsheet': 'CSV Spreadsheet',
        'generate_download': 'Generate & Download',
        'category_breakdown': 'Category Breakdown',
        'status': 'Status',
        'spent': 'Spent',
        'budget': 'Budget',
        'safe': 'Safe',
        'over_limit': 'Over Limit',
        'no_limit_set': 'No Limit Set',
        'welcome': 'Welcome',
        'logout': 'Logout'
    },
    'es': {
        'dashboard': 'Tablero',
        'transactions': 'Transacciones',
        'reports': 'Informes',
        'settings': 'Ajustes',
        'total_balance': 'Balance Total',
        'monthly_income': 'Ingresos Mensuales',
        'monthly_expenses': 'Gastos Mensuales',
        'monthly_savings': 'Ahorros Mensuales',
        'savings_rate': 'Tasa de Ahorro',
        'budget_consumptions': 'Consumo de Presupuesto',
        'highest_spending_category': 'Categoría de Mayor Gasto',
        'spent_breakdown': 'Desglose de Gastos',
        'spending_trends': 'Tendencias de Gasto',
        'upcoming_bills': 'Próximas Facturas',
        'recent_transactions': 'Transacciones Recientes',
        'add_transaction': 'Agregar Transacción',
        'set_budget': 'Definir Presupuesto',
        'add_bill': 'Agregar Factura',
        'export_options': 'Opciones de Exportación',
        'select_month': 'Seleccionar Mes',
        'export_format': 'Formato de Exportación',
        'pdf_document': 'Documento PDF',
        'csv_spreadsheet': 'Hoja de Cálculo CSV',
        'generate_download': 'Generar y Descargar',
        'category_breakdown': 'Desglose por Categoría',
        'status': 'Estado',
        'spent': 'Gastado',
        'budget': 'Presupuesto',
        'safe': 'Seguro',
        'over_limit': 'Límite Excedido',
        'no_limit_set': 'Sin Límite',
        'welcome': 'Bienvenido',
        'logout': 'Cerrar Sesión'
    },
    'hi': {
        'dashboard': 'डैशबोर्ड',
        'transactions': 'लेन-देन',
        'reports': 'रिपोर्ट',
        'settings': 'सेटिंग्स',
        'total_balance': 'कुल शेष राशि',
        'monthly_income': 'मासिक आय',
        'monthly_expenses': 'मासिक खर्च',
        'monthly_savings': 'मासिक बचत',
        'savings_rate': 'बचत दर',
        'budget_consumptions': 'बजट की खपत',
        'highest_spending_category': 'उच्चतम व्यय श्रेणी',
        'spent_breakdown': 'खर्च विवरण',
        'spending_trends': 'व्यय प्रवृत्तियां',
        'upcoming_bills': 'आने वाले बिल',
        'recent_transactions': 'हाल ही के लेनदेन',
        'add_transaction': 'लेनदेन जोड़ें',
        'set_budget': 'बजट निर्धारित करें',
        'add_bill': 'बिल जोड़ें',
        'export_options': 'निर्यात विकल्प',
        'select_month': 'महीना चुनें',
        'export_format': 'निर्यात प्रारूप',
        'pdf_document': 'पीडीएफ दस्तावेज़',
        'csv_spreadsheet': 'सीएसवी स्प्रेडशीट',
        'generate_download': 'बनाएं और डाउनलोड करें',
        'category_breakdown': 'श्रेणी विवरण',
        'status': 'स्थिति',
        'spent': 'खर्च किया',
        'budget': 'बजट',
        'safe': 'सुरक्षित',
        'over_limit': 'सीमा पार',
        'no_limit_set': 'कोई सीमा नहीं',
        'welcome': 'स्वागत हे',
        'logout': 'लॉग आउट'
    }
}

@app.context_processor
def inject_translations():
    lang = 'en'
    if current_user and current_user.is_authenticated:
        lang = current_user.language or 'en'
    def translate(key):
        return TRANSLATIONS.get(lang, TRANSLATIONS['en']).get(key, key)
    return dict(t=translate)

def get_notifications():
    notifications = []
    if not current_user or not current_user.is_authenticated:
        return notifications
        
    today = date.today()
    current_month_str = today.strftime("%Y-%m")
    symbol = current_user.currency or '₹'
    
    # 1. Daily expense reminder
    if current_user.daily_reminder:
        has_tx_today = Transaction.query.filter_by(user_id=current_user.id, date=today).first() is not None
        if not has_tx_today:
            notifications.append({
                'type': 'warning',
                'text': 'Daily Reminder: You haven\'t logged any expenses today yet!'
            })
            
    # 2. Weekly spending summary
    if current_user.weekly_summary:
        seven_days_ago = today - timedelta(days=7)
        weekly_spent = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'expense',
            Transaction.date >= seven_days_ago
        ).scalar() or 0.0
        if weekly_spent > 0:
            notifications.append({
                'type': 'info',
                'text': f'Weekly Summary: You spent {symbol}{weekly_spent:,.2f} in the past 7 days.'
            })
            
    # 3. Budget threshold warnings
    budgets = Budget.query.filter_by(user_id=current_user.id, month=current_month_str).all()
    for b in budgets:
        if b.category == 'all':
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                extract('year', Transaction.date) == today.year,
                extract('month', Transaction.date) == today.month
            ).scalar() or 0.0
        else:
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                Transaction.category == b.category,
                extract('year', Transaction.date) == today.year,
                extract('month', Transaction.date) == today.month
            ).scalar() or 0.0
            
        threshold_val = current_user.budget_warning_threshold or 90.0
        percent = (spent / b.amount * 100.0) if b.amount > 0 else 0.0
        if percent >= threshold_val:
            cat_name = "Overall" if b.category == 'all' else b.category.capitalize()
            notifications.append({
                'type': 'danger',
                'text': f'Budget Alert: "{cat_name}" is at {percent:.1f}% consumed ({symbol}{spent:,.2f} / {symbol}{b.amount:,.2f})!'
            })
            
    # 4. Savings Milestones alerts
    if current_user.savings_milestones:
        # Lifetime total net balance
        total_income = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id, Transaction.type == 'income'
        ).scalar() or 0.0
        total_expense = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id, Transaction.type == 'expense'
        ).scalar() or 0.0
        balance = total_income - total_expense
        
        milestones = [100000, 50000, 25000, 10000, 5000]
        for val in milestones:
            if balance >= val:
                notifications.append({
                    'type': 'success',
                    'text': f'Milestone Unlocked: Your lifetime net balance crossed {symbol}{val:,}!'
                })
                break
                
    return notifications

@app.context_processor
def inject_notifications_list():
    return dict(notifications=get_notifications())

def get_user_categories():
    defaults = ["Food", "Travel", "Shopping", "Bills", "Entertainment", "Other"]
    if current_user and current_user.is_authenticated:
        customs = [c.name for c in current_user.custom_categories]
        return list(dict.fromkeys(defaults + customs))
    return defaults

# --- Authentication Routes ---
@app.route('/register', methods=['GET', 'POST'])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    if request.method == 'POST':
        username = request.form.get('username')
        email = request.form.get('email')
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')

        if password != confirm_password:
            flash("Passwords do not match.", "error")
            return redirect(url_for('register'))

        existing_user = User.query.filter((User.username == username) | (User.email == email)).first()
        if existing_user:
            flash("Username or email already exists.", "error")
            return redirect(url_for('register'))

        new_user = User(username=username, email=email)
        new_user.set_password(password)
        db.session.add(new_user)
        db.session.commit()

        flash("Registration successful. Please login.", "success")
        return redirect(url_for('login'))

    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')

        user = User.query.filter_by(username=username).first()
        if user and user.check_password(password):
            login_user(user)
            flash("Successfully logged in.", "success")
            return redirect(url_for('dashboard'))
        else:
            flash("Invalid username or password.", "error")
            return redirect(url_for('login'))

    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    flash("Successfully logged out.", "success")
    return redirect(url_for('login'))

# --- Application Routes ---
@app.route('/')
@app.route('/dashboard')
@login_required
def dashboard():
    today = date.today()
    current_month_str = today.strftime("%Y-%m")
    
    # 1. Total Income (this month)
    total_income = db.session.query(db.func.sum(Transaction.amount)).filter(
        Transaction.user_id == current_user.id,
        Transaction.type == 'income',
        extract('year', Transaction.date) == today.year,
        extract('month', Transaction.date) == today.month
    ).scalar() or 0.0

    # 2. Total Expense (this month)
    total_expense = db.session.query(db.func.sum(Transaction.amount)).filter(
        Transaction.user_id == current_user.id,
        Transaction.type == 'expense',
        extract('year', Transaction.date) == today.year,
        extract('month', Transaction.date) == today.month
    ).scalar() or 0.0

    # 3. Remaining Balance (all-time Income - Expense)
    all_time_income = db.session.query(db.func.sum(Transaction.amount)).filter(
        Transaction.user_id == current_user.id,
        Transaction.type == 'income'
    ).scalar() or 0.0

    all_time_expense = db.session.query(db.func.sum(Transaction.amount)).filter(
        Transaction.user_id == current_user.id,
        Transaction.type == 'expense'
    ).scalar() or 0.0

    remaining_balance = all_time_income - all_time_expense

    # 4. Savings (Monthly)
    monthly_savings = total_income - total_expense
    savings_rate = (monthly_savings / total_income * 100) if total_income > 0 else 0.0

    # 5. Recent Transactions (last 5)
    recent_transactions = Transaction.query.filter_by(user_id=current_user.id)\
        .order_by(Transaction.date.desc(), Transaction.id.desc()).limit(5).all()

    # 6. Expense by Category (for current month pie chart)
    category_data = {}
    categories_query = db.session.query(
        Transaction.category, db.func.sum(Transaction.amount)
    ).filter(
        Transaction.user_id == current_user.id,
        Transaction.type == 'expense',
        extract('year', Transaction.date) == today.year,
        extract('month', Transaction.date) == today.month
    ).group_by(Transaction.category).all()
    
    for cat, amt in categories_query:
        category_data[cat] = amt

    # 7. Top Spending Category
    top_category = "N/A"
    top_category_amount = 0.0
    if category_data:
        top_category = max(category_data, key=category_data.get)
        top_category_amount = category_data[top_category]

    # 8. Monthly Spending (for bar chart - last 6 months)
    monthly_data = {}
    for i in range(5, -1, -1):
        m = today.month - i
        y = today.year
        while m <= 0:
            m += 12
            y -= 1
        
        month_name = datetime(y, m, 1).strftime("%b %Y")
        
        expense_sum = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'expense',
            extract('year', Transaction.date) == y,
            extract('month', Transaction.date) == m
        ).scalar() or 0.0
        
        monthly_data[month_name] = expense_sum

    # 9. Budgets and alerts
    active_budgets = []
    budget_alerts = []
    
    budgets = Budget.query.filter_by(user_id=current_user.id, month=current_month_str).all()
    for b in budgets:
        if b.category == 'all':
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                extract('year', Transaction.date) == today.year,
                extract('month', Transaction.date) == today.month
            ).scalar() or 0.0
            
            percent = min(int(spent / b.amount * 100), 100) if b.amount > 0 else 0
            active_budgets.append({
                'category': 'Overall Budget',
                'limit': b.amount,
                'spent': spent,
                'percent': percent,
                'exceeded': spent > b.amount
            })
            if spent > b.amount:
                budget_alerts.append(f"Overall monthly spending (₹{spent:.2f}) has exceeded your limit of ₹{b.amount:.2f}!")
        else:
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                Transaction.category == b.category,
                extract('year', Transaction.date) == today.year,
                extract('month', Transaction.date) == today.month
            ).scalar() or 0.0
            
            percent = min(int(spent / b.amount * 100), 100) if b.amount > 0 else 0
            active_budgets.append({
                'category': b.category,
                'limit': b.amount,
                'spent': spent,
                'percent': percent,
                'exceeded': spent > b.amount
            })
            if spent > b.amount:
                budget_alerts.append(f"Spending in category '{b.category}' (₹{spent:.2f}) has exceeded your budget limit of ₹{b.amount:.2f}!")

    # 10. Financial Health Score Algorithm
    financial_score = 50  # base
    
    # Savings component (max +30 points)
    if savings_rate > 0:
        financial_score += min(savings_rate * 0.6, 30)
    elif savings_rate < 0:
        financial_score += max(savings_rate * 0.4, -20)
        
    # Budget overruns check (max -30 points)
    overruns = sum(1 for b in active_budgets if b['exceeded'])
    financial_score -= min(overruns * 10, 30)
    
    # Budget adherence award
    if len(active_budgets) > 0 and overruns == 0:
        financial_score += 15
        
    # Activity check
    if len(recent_transactions) > 0:
        financial_score += 5
        
    financial_score = max(0, min(100, int(financial_score)))
    
    if financial_score >= 80:
        score_tier = "Excellent"
        score_color = "success"
    elif financial_score >= 60:
        score_tier = "Good"
        score_color = "warning"
    else:
        score_tier = "Needs Attention"
        score_color = "danger"

    # 11. Upcoming Bills
    upcoming_bills = UpcomingBill.query.filter_by(user_id=current_user.id, is_paid=False)\
        .order_by(UpcomingBill.due_date.asc()).all()

    return render_template(
        'dashboard.html',
        total_income=total_income,
        total_expense=total_expense,
        remaining_balance=remaining_balance,
        monthly_savings=monthly_savings,
        savings_rate=savings_rate,
        recent_transactions=recent_transactions,
        category_data=category_data,
        top_category=top_category,
        top_category_amount=top_category_amount,
        monthly_data=monthly_data,
        active_budgets=active_budgets,
        budget_alerts=budget_alerts,
        financial_score=financial_score,
        score_tier=score_tier,
        score_color=score_color,
        upcoming_bills=upcoming_bills,
        current_month=current_month_str,
        categories=get_user_categories()
    )

# --- Bill Routing ---
@app.route('/bill/add', methods=['POST'])
@login_required
def add_bill():
    title = request.form.get('title')
    amount_str = request.form.get('amount')
    due_date_str = request.form.get('due_date')
    category = request.form.get('category')

    try:
        amount = float(amount_str)
        if amount <= 0:
            raise ValueError
    except (ValueError, TypeError):
        flash("Invalid bill amount.", "error")
        return redirect(url_for('dashboard'))

    try:
        due_date = datetime.strptime(due_date_str, '%Y-%m-%d').date()
    except (ValueError, TypeError):
        flash("Invalid due date format.", "error")
        return redirect(url_for('dashboard'))

    new_bill = UpcomingBill(
        user_id=current_user.id,
        title=title,
        amount=amount,
        due_date=due_date,
        category=category,
        is_paid=False
    )
    db.session.add(new_bill)
    db.session.commit()
    flash(f"Upcoming bill '{title}' successfully added.", "success")
    return redirect(url_for('dashboard'))

@app.route('/bill/pay/<int:id>', methods=['POST'])
@login_required
def pay_bill(id):
    bill = UpcomingBill.query.get_or_404(id)
    if bill.user_id != current_user.id:
        flash("Unauthorized action.", "error")
        return redirect(url_for('dashboard'))

    # Mark paid and log as expense
    new_t = Transaction(
        user_id=current_user.id,
        type='expense',
        category=bill.category,
        amount=bill.amount,
        description=f"Paid Bill: {bill.title}",
        date=date.today()
    )
    bill.is_paid = True
    db.session.add(new_t)
    db.session.commit()
    
    flash(f"Bill '{bill.title}' successfully marked as paid and logged as a transaction.", "success")
    return redirect(url_for('dashboard'))

# --- Transactions Log Routing ---
@app.route('/transactions')
@login_required
def transactions():
    search_query = request.args.get('search', '')
    type_filter = request.args.get('type', '')
    category_filter = request.args.get('category', '')
    start_date = request.args.get('start_date', '')
    end_date = request.args.get('end_date', '')

    query = Transaction.query.filter_by(user_id=current_user.id)

    if search_query:
        query = query.filter(Transaction.description.ilike(f"%{search_query}%"))
    if type_filter:
        query = query.filter(Transaction.type == type_filter)
    if category_filter:
        query = query.filter(Transaction.category == category_filter)
    if start_date:
        try:
            s_date = datetime.strptime(start_date, '%Y-%m-%d').date()
            query = query.filter(Transaction.date >= s_date)
        except ValueError:
            pass
    if end_date:
        try:
            e_date = datetime.strptime(end_date, '%Y-%m-%d').date()
            query = query.filter(Transaction.date <= e_date)
        except ValueError:
            pass

    results = query.order_by(Transaction.date.desc(), Transaction.id.desc()).all()

    return render_template(
        'transactions.html',
        transactions=results,
        search_query=search_query,
        type_filter=type_filter,
        category_filter=category_filter,
        start_date=start_date,
        end_date=end_date,
        categories=get_user_categories()
    )

@app.route('/transaction/add', methods=['GET', 'POST'])
@login_required
def add_transaction():
    if request.method == 'POST':
        type_val = request.form.get('type')
        amount_str = request.form.get('amount')
        category = request.form.get('category')
        date_str = request.form.get('date')
        description = request.form.get('description')

        try:
            amount = float(amount_str)
        except (ValueError, TypeError):
            flash("Invalid amount.", "error")
            return redirect(url_for('add_transaction'))

        try:
            date_val = datetime.strptime(date_str, '%Y-%m-%d').date()
        except (ValueError, TypeError):
            flash("Invalid date format.", "error")
            return redirect(url_for('add_transaction'))

        receipt_filename = None
        file = request.files.get('receipt')
        if file and file.filename != '':
            if allowed_file(file.filename):
                filename = secure_filename(file.filename)
                receipt_filename = f"{uuid.uuid4().hex}_{filename}"
                file.save(os.path.join(UPLOAD_FOLDER, receipt_filename))
            else:
                flash("Unsupported receipt format. Only images are allowed.", "error")

        new_t = Transaction(
            user_id=current_user.id,
            type=type_val,
            category=category,
            amount=amount,
            description=description,
            date=date_val,
            receipt_image=receipt_filename
        )
        db.session.add(new_t)
        db.session.commit()

        # Check category budget alert instantly
        month_str = date_val.strftime("%Y-%m")
        budget = Budget.query.filter_by(user_id=current_user.id, category=category, month=month_str).first()
        if budget:
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                Transaction.category == category,
                extract('year', Transaction.date) == date_val.year,
                extract('month', Transaction.date) == date_val.month
            ).scalar() or 0.0
            if spent > budget.amount:
                flash(f"Budget Exceeded! You have spent ₹{spent:.2f} on {category} against budget of ₹{budget.amount:.2f}.", "error")

        flash("Transaction added successfully.", "success")
        return redirect(url_for('dashboard'))

    today_str = date.today().strftime('%Y-%m-%d')
    return render_template('add_transaction.html', today=today_str, transaction=None, categories=get_user_categories())

@app.route('/transaction/edit/<int:id>', methods=['GET', 'POST'])
@login_required
def edit_transaction(id):
    t = Transaction.query.get_or_404(id)
    if t.user_id != current_user.id:
        flash("Unauthorized action.", "error")
        return redirect(url_for('dashboard'))

    if request.method == 'POST':
        type_val = request.form.get('type')
        amount_str = request.form.get('amount')
        category = request.form.get('category')
        date_str = request.form.get('date')
        description = request.form.get('description')

        try:
            amount = float(amount_str)
        except (ValueError, TypeError):
            flash("Invalid amount.", "error")
            return redirect(url_for('edit_transaction', id=id))

        try:
            date_val = datetime.strptime(date_str, '%Y-%m-%d').date()
        except (ValueError, TypeError):
            flash("Invalid date format.", "error")
            return redirect(url_for('edit_transaction', id=id))

        t.type = type_val
        t.category = category
        t.amount = amount
        t.date = date_val
        t.description = description

        file = request.files.get('receipt')
        if file and file.filename != '':
            if allowed_file(file.filename):
                # Clean old file if exists
                if t.receipt_image:
                    old_path = os.path.join(UPLOAD_FOLDER, t.receipt_image)
                    if os.path.exists(old_path):
                        try:
                            os.remove(old_path)
                        except OSError:
                            pass
                
                filename = secure_filename(file.filename)
                t.receipt_image = f"{uuid.uuid4().hex}_{filename}"
                file.save(os.path.join(UPLOAD_FOLDER, t.receipt_image))
            else:
                flash("Unsupported receipt format. Only images are allowed.", "error")

        db.session.commit()
        flash("Transaction updated successfully.", "success")
        return redirect(url_for('transactions'))

    today_str = date.today().strftime('%Y-%m-%d')
    return render_template('add_transaction.html', transaction=t, today=today_str, categories=get_user_categories())

@app.route('/transaction/delete/<int:id>', methods=['POST'])
@login_required
def delete_transaction(id):
    t = Transaction.query.get_or_404(id)
    if t.user_id != current_user.id:
        flash("Unauthorized action.", "error")
        return redirect(url_for('dashboard'))

    if t.receipt_image:
        filepath = os.path.join(UPLOAD_FOLDER, t.receipt_image)
        if os.path.exists(filepath):
            try:
                os.remove(filepath)
            except OSError:
                pass

    db.session.delete(t)
    db.session.commit()
    flash("Transaction deleted successfully.", "success")
    return redirect(request.referrer or url_for('dashboard'))

# --- Budget Setting Routing ---
@app.route('/budget/set', methods=['POST'])
@login_required
def set_budget():
    category = request.form.get('category')
    amount_str = request.form.get('amount')
    month = request.form.get('month')  # Format: YYYY-MM

    try:
        amount = float(amount_str)
    except (ValueError, TypeError):
        flash("Invalid budget amount.", "error")
        return redirect(url_for('dashboard'))

    # Check if budget limit exists
    existing = Budget.query.filter_by(
        user_id=current_user.id,
        category=category,
        month=month
    ).first()

    if existing:
        existing.amount = amount
        flash(f"Budget for '{category}' successfully updated to ₹{amount:.2f}.", "success")
    else:
        new_b = Budget(
            user_id=current_user.id,
            category=category,
            amount=amount,
            month=month
        )
        db.session.add(new_b)
        flash(f"Budget for '{category}' successfully created for ₹{amount:.2f}.", "success")

    db.session.commit()
    return redirect(request.referrer or url_for('dashboard'))

# --- User Settings & Personalization ---
@app.route('/settings', methods=['GET', 'POST'])
@login_required
def settings():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        email = request.form.get('email', '').strip()
        currency = request.form.get('currency', '₹')
        language = request.form.get('language', 'en')
        theme_accent = request.form.get('theme_accent', 'indigo')
        daily_reminder = 'daily_reminder' in request.form
        weekly_summary = 'weekly_summary' in request.form
        savings_milestones = 'savings_milestones' in request.form
        
        try:
            budget_warning_threshold = float(request.form.get('budget_warning_threshold', 90.0))
        except (ValueError, TypeError):
            budget_warning_threshold = 90.0

        if not username or not email:
            flash("Username and Email are required.", "error")
            return redirect(url_for('settings'))

        existing_user = User.query.filter(User.username == username, User.id != current_user.id).first()
        if existing_user:
            flash("Username already taken.", "error")
            return redirect(url_for('settings'))

        existing_email = User.query.filter(User.email == email, User.id != current_user.id).first()
        if existing_email:
            flash("Email already registered.", "error")
            return redirect(url_for('settings'))

        current_user.username = username
        current_user.email = email
        current_user.currency = currency
        current_user.language = language
        current_user.theme_accent = theme_accent
        current_user.daily_reminder = daily_reminder
        current_user.weekly_summary = weekly_summary
        current_user.savings_milestones = savings_milestones
        current_user.budget_warning_threshold = budget_warning_threshold

        current_password = request.form.get('current_password', '')
        new_password = request.form.get('new_password', '')
        if current_password and new_password:
            if current_user.check_password(current_password):
                current_user.set_password(new_password)
                flash("Password updated successfully.", "success")
            else:
                flash("Incorrect current password.", "error")
                return redirect(url_for('settings'))

        file = request.files.get('profile_pic')
        if file and file.filename != '':
            if allowed_file(file.filename):
                if current_user.profile_pic:
                    old_path = os.path.join(PROFILE_PICS_FOLDER, current_user.profile_pic)
                    if os.path.exists(old_path):
                        try:
                            os.remove(old_path)
                        except OSError:
                            pass
                
                filename = secure_filename(file.filename)
                unique_filename = f"{uuid.uuid4().hex}_{filename}"
                file.save(os.path.join(PROFILE_PICS_FOLDER, unique_filename))
                current_user.profile_pic = unique_filename
                flash("Profile picture updated.", "success")
            else:
                flash("Invalid profile picture format.", "error")
                return redirect(url_for('settings'))

        db.session.commit()
        flash("Settings updated successfully.", "success")
        return redirect(url_for('settings'))

    return render_template('settings.html')

@app.route('/settings/category/add', methods=['POST'])
@login_required
def add_custom_category():
    name = request.form.get('name', '').strip()
    if not name:
        flash("Category name cannot be empty.", "error")
        return redirect(url_for('settings'))

    existing = CustomCategory.query.filter_by(user_id=current_user.id, name=name).first()
    if existing or name.capitalize() in ["Food", "Travel", "Shopping", "Bills", "Entertainment", "Other"]:
        flash(f"Category '{name}' already exists.", "error")
        return redirect(url_for('settings'))

    new_cat = CustomCategory(user_id=current_user.id, name=name)
    db.session.add(new_cat)
    db.session.commit()
    flash(f"Custom category '{name}' added successfully.", "success")
    return redirect(url_for('settings'))

@app.route('/settings/category/delete/<int:id>', methods=['POST'])
@login_required
def delete_custom_category(id):
    cat = CustomCategory.query.get_or_404(id)
    if cat.user_id != current_user.id:
        flash("Unauthorized action.", "error")
        return redirect(url_for('settings'))

    db.session.delete(cat)
    db.session.commit()
    flash(f"Custom category deleted successfully.", "success")
    return redirect(url_for('settings'))

# --- Reports & Analytical breakdown Routing ---
@app.route('/reports')
@login_required
def reports():
    view_month = request.args.get('view_month', '')
    if not view_month:
        view_month = date.today().strftime('%Y-%m')

    # Convert YYYY-MM into datetime objects for queries
    try:
        y, m = map(int, view_month.split('-'))
        selected_date = datetime(y, m, 1)
    except ValueError:
        today = date.today()
        y, m = today.year, today.month
        selected_date = datetime(y, m, 1)
        view_month = today.strftime('%Y-%m')

    selected_month_name = selected_date.strftime("%B %Y")

    # Load category breakdowns
    # We want a list of category analysis: {category, spent, budget}
    category_analysis = []
    
    # Pre-defined categories
    categories = ["Food", "Travel", "Shopping", "Bills", "Entertainment", "Other"]
    
    for cat in categories:
        spent = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'expense',
            Transaction.category == cat,
            extract('year', Transaction.date) == y,
            extract('month', Transaction.date) == m
        ).scalar() or 0.0

        budget = Budget.query.filter_by(
            user_id=current_user.id,
            category=cat,
            month=view_month
        ).first()
        budget_limit = budget.amount if budget else 0.0

        if spent > 0 or budget_limit > 0:
            category_analysis.append({
                'category': cat,
                'spent': spent,
                'budget': budget_limit
            })

    return render_template(
        'reports.html',
        selected_month=view_month,
        selected_month_name=selected_month_name,
        category_analysis=category_analysis
    )

# --- CSV & PDF Exports Router ---
@app.route('/export')
@login_required
def export_report():
    month_str = request.args.get('month', '')
    export_format = request.args.get('format', 'pdf')

    if not month_str:
        month_str = date.today().strftime('%Y-%m')

    return redirect(url_for('export_report_path', month_str=month_str, export_format=export_format))

@app.route('/export/<month_str>.<export_format>')
@login_required
def export_report_path(month_str, export_format):
    try:
        y, m = map(int, month_str.split('-'))
        selected_month_name = datetime(y, m, 1).strftime("%B %Y")
    except ValueError:
        flash("Invalid month selected for export.", "error")
        return redirect(url_for('reports'))

    # Load transactions for the month
    transactions = Transaction.query.filter(
        Transaction.user_id == current_user.id,
        extract('year', Transaction.date) == y,
        extract('month', Transaction.date) == m
    ).order_by(Transaction.date.asc()).all()

    if export_format == 'csv':
        # CSV Export
        si = io.StringIO()
        cw = csv.writer(si)
        # Headers
        cw.writerow(["Date", "Description", "Category", "Type", "Amount"])
        for t in transactions:
            cw.writerow([
                t.date.strftime('%Y-%m-%d'),
                t.description or 'No Description',
                t.category,
                t.type.capitalize(),
                f"{t.amount:.2f}"
            ])
        
        output = make_response_csv(si.getvalue(), f"Expense_Report_{month_str}.csv")
        return output

    elif export_format == 'pptx':
        total_income = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'income',
            extract('year', Transaction.date) == y,
            extract('month', Transaction.date) == m
        ).scalar() or 0.0

        total_expense = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'expense',
            extract('year', Transaction.date) == y,
            extract('month', Transaction.date) == m
        ).scalar() or 0.0

        category_analysis = []
        categories = ["Food", "Travel", "Shopping", "Bills", "Entertainment", "Other"]
        for cat in categories:
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                Transaction.category == cat,
                extract('year', Transaction.date) == y,
                extract('month', Transaction.date) == m
            ).scalar() or 0.0

            budget = Budget.query.filter_by(
                user_id=current_user.id,
                category=cat,
                month=month_str
            ).first()
            budget_limit = budget.amount if budget else 0.0

            if spent > 0 or budget_limit > 0:
                category_analysis.append({
                    'category': cat,
                    'spent': spent,
                    'budget': budget_limit
                })

        try:
            pptx_bytes = generate_pptx_data(
                user=current_user,
                month_str=selected_month_name,
                transactions=transactions,
                total_income=total_income,
                total_expense=total_expense,
                category_analysis=category_analysis
            )
            response = send_file(
                io.BytesIO(pptx_bytes),
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f"Expense_Report_{month_str}.pptx"
            )
            response.headers['Content-Disposition'] = f'attachment; filename="Expense_Report_{month_str}.pptx"'
            return response
        except Exception as e:
            flash(f"Error generating PPTX: {str(e)}", "error")
            return redirect(url_for('reports'))

    else:
        # PDF Export
        total_income = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'income',
            extract('year', Transaction.date) == y,
            extract('month', Transaction.date) == m
        ).scalar() or 0.0

        total_expense = db.session.query(db.func.sum(Transaction.amount)).filter(
            Transaction.user_id == current_user.id,
            Transaction.type == 'expense',
            extract('year', Transaction.date) == y,
            extract('month', Transaction.date) == m
        ).scalar() or 0.0

        category_analysis = []
        categories = ["Food", "Travel", "Shopping", "Bills", "Entertainment", "Other"]
        for cat in categories:
            spent = db.session.query(db.func.sum(Transaction.amount)).filter(
                Transaction.user_id == current_user.id,
                Transaction.type == 'expense',
                Transaction.category == cat,
                extract('year', Transaction.date) == y,
                extract('month', Transaction.date) == m
            ).scalar() or 0.0

            budget = Budget.query.filter_by(
                user_id=current_user.id,
                category=cat,
                month=month_str
            ).first()
            budget_limit = budget.amount if budget else 0.0

            if spent > 0 or budget_limit > 0:
                category_analysis.append({
                    'category': cat,
                    'spent': spent,
                    'budget': budget_limit
                })

        try:
            pdf_bytes = generate_pdf_data(
                user=current_user,
                month_str=selected_month_name,
                transactions=transactions,
                total_income=total_income,
                total_expense=total_expense,
                category_analysis=category_analysis
            )
            response = send_file(
                io.BytesIO(pdf_bytes),
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f"Expense_Report_{month_str}.pdf"
            )
            response.headers['Content-Disposition'] = f'attachment; filename="Expense_Report_{month_str}.pdf"'
            return response
        except Exception as e:
            flash(f"Error generating PDF: {str(e)}", "error")
            return redirect(url_for('reports'))

def make_response_csv(csv_text, filename):
    return Response(
        csv_text,
        mimetype="text/csv",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'}
    )

def generate_pdf_data(user, month_str, transactions, total_income, total_expense, category_analysis):
    from html import escape
    symbol = user.currency or '₹'
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    story = []
    
    styles = getSampleStyleSheet()
    
    # Custom report styles
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=22,
        leading=26,
        textColor=colors.HexColor('#4f46e5'),
        spaceAfter=5
    )
    subtitle_style = ParagraphStyle(
        'SubtitleStyle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#64748b'),
        spaceAfter=25
    )
    h2_style = ParagraphStyle(
        'H2Style',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=14,
        leading=18,
        textColor=colors.HexColor('#0f172a'),
        spaceBefore=15,
        spaceAfter=10
    )
    body_style = ParagraphStyle(
        'BodyStyle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        leading=12,
        textColor=colors.HexColor('#0f172a')
    )
    body_bold = ParagraphStyle(
        'BodyBold',
        parent=body_style,
        fontName='Helvetica-Bold'
    )
    
    # 1. Header
    story.append(Paragraph(f"Expense Tracker Report", title_style))
    story.append(Paragraph(f"Monthly financial statement for {month_str} | Generated by {escape(user.username)}", subtitle_style))
    
    # 2. Key Metrics Card
    net_savings = total_income - total_expense
    savings_color = "#059669" if net_savings >= 0 else "#dc2626"
    
    summary_data = [
        [
            Paragraph("<b>Total Income</b>", body_style),
            Paragraph("<b>Total Expense</b>", body_style),
            Paragraph("<b>Net Savings</b>", body_style)
        ],
        [
            Paragraph(f"<font color='#059669'>+{symbol}{total_income:.2f}</font>", body_bold),
            Paragraph(f"<font color='#dc2626'>-{symbol}{total_expense:.2f}</font>", body_bold),
            Paragraph(f"<font color='{savings_color}'>{symbol}{net_savings:.2f}</font>", body_bold)
        ]
    ]
    summary_table = Table(summary_data, colWidths=[175, 175, 175])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#f8fafc')),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('PADDING', (0,0), (-1,-1), 12),
        ('LINEBELOW', (0,0), (-1,0), 1, colors.HexColor('#e2e8f0')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#e2e8f0')),
    ]))
    story.append(summary_table)
    story.append(Spacer(1, 15))
    
    # 3. Category Summary Table
    story.append(Paragraph("Category Spending & Budgets", h2_style))
    cat_data = [[
        Paragraph("<b>Category</b>", body_style),
        Paragraph("<b>Spent</b>", body_style),
        Paragraph("<b>Budget Limit</b>", body_style),
        Paragraph("<b>Status</b>", body_style)
    ]]
    
    for item in category_analysis:
        status_text = "Safe"
        status_color = "#059669"
        if item['budget'] > 0:
            if item['spent'] > item['budget']:
                status_text = "Over Limit"
                status_color = "#dc2626"
        else:
            status_text = "No Limit"
            status_color = "#64748b"
            
        budget_text = f"{symbol}{item['budget']:.2f}" if item['budget'] > 0 else "—"
        cat_data.append([
            Paragraph(escape(item['category'].capitalize()), body_style),
            Paragraph(f"{symbol}{item['spent']:.2f}", body_style),
            Paragraph(budget_text, body_style),
            Paragraph(f"<font color='{status_color}'><b>{status_text}</b></font>", body_style)
        ])
        
    if len(category_analysis) == 0:
        cat_data.append([Paragraph("No spending recorded in this month", body_style), "", "", ""])
    
    cat_table = Table(cat_data, colWidths=[130, 120, 120, 155])
    cat_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f1f5f9')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('PADDING', (0,0), (-1,-1), 8),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
    ]))
    story.append(cat_table)
    story.append(Spacer(1, 15))
    
    # 4. Transactions List Table
    story.append(Paragraph("Transactions Details", h2_style))
    tx_data = [[
        Paragraph("<b>Date</b>", body_style),
        Paragraph("<b>Description</b>", body_style),
        Paragraph("<b>Category</b>", body_style),
        Paragraph("<b>Type</b>", body_style),
        Paragraph("<b>Amount</b>", body_style)
    ]]
    
    for t in transactions:
        amt_prefix = "+" if t.type == 'income' else "-"
        amt_color = "#059669" if t.type == 'income' else "#dc2626"
        tx_data.append([
            Paragraph(t.date.strftime('%b %d, %Y'), body_style),
            Paragraph(escape(t.description or 'No Description'), body_style),
            Paragraph(escape(t.category), body_style),
            Paragraph(t.type.capitalize(), body_style),
            Paragraph(f"<font color='{amt_color}'><b>{amt_prefix}{symbol}{t.amount:.2f}</b></font>", body_style)
        ])
        
    if len(transactions) == 0:
        tx_data.append([Paragraph("No transactions recorded in this month", body_style), "", "", "", ""])
        
    tx_table = Table(tx_data, colWidths=[80, 165, 100, 70, 110])
    tx_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#f1f5f9')),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('PADDING', (0,0), (-1,-1), 8),
        ('LINEBELOW', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
    ]))
    story.append(tx_table)
    
    doc.build(story)
    pdf_data = buffer.getvalue()
    buffer.close()
    return pdf_data

def generate_pptx_data(user, month_str, transactions, total_income, total_expense, category_analysis):
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    DARK_BG = RGBColor(9, 13, 22)
    EMERALD = RGBColor(16, 185, 129)
    INDIGO = RGBColor(99, 102, 241)
    WHITE = RGBColor(248, 250, 252)
    SLATE = RGBColor(148, 163, 184)
    CARD_BG = RGBColor(17, 24, 39)
    
    def set_slide_background(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    slide_layout = prs.slide_layouts[6]
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, DARK_BG)
    
    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Monthly Financial Report"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = EMERALD
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"Statement for {month_str} | Generated by {user.username}"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(20)
    p2.font.color.rgb = INDIGO
    p2.space_before = Pt(10)
    p2.alignment = PP_ALIGN.CENTER

    def add_slide_header(slide, title_text):
        catBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        cat_tf = catBox.text_frame
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = f"MONTHLY STATEMENT — {user.username.upper()}"
        cat_p.font.name = 'Trebuchet MS'
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = INDIGO
        
        titleBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = titleBox.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = 'Trebuchet MS'
        title_p.font.size = Pt(26)
        title_p.font.bold = True
        title_p.font.color.rgb = WHITE
        
        connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.04))
        connector.fill.solid()
        connector.fill.fore_color.rgb = EMERALD
        connector.line.fill.background()

    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, DARK_BG)
    add_slide_header(slide2, "Summary Statistics")
    
    net_savings = total_income - total_expense
    symbol = user.currency or '₹'
    
    def add_summary_card(slide, left, top, width, height, title, value_str, val_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = INDIGO
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.color.rgb = SLATE
        p.space_after = Pt(12)
        
        p2 = tf.add_paragraph()
        p2.text = value_str
        p2.font.name = 'Trebuchet MS'
        p2.font.size = Pt(30)
        p2.font.bold = True
        p2.font.color.rgb = val_color
        
    add_summary_card(slide2, 0.8, 2.2, 3.6, 3.5, "Total Income", f"+{symbol}{total_income:,.2f}", EMERALD)
    add_summary_card(slide2, 4.8, 2.2, 3.6, 3.5, "Total Expenses", f"-{symbol}{total_expense:,.2f}", RGBColor(239, 68, 68))
    add_summary_card(slide2, 8.8, 2.2, 3.7, 3.5, "Net Savings", f"{symbol}{net_savings:,.2f}", EMERALD if net_savings >= 0 else RGBColor(239, 68, 68))

    # Slide 3: Category Spent Breakdown Table
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3, DARK_BG)
    add_slide_header(slide3, "Category Spending & Budgets")
    
    rows = len(category_analysis) + 1
    cols = 4
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(11.7)
    height = Inches(0.5 * rows)
    
    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Headers
    headers = ["Category", "Spent Amount", "Budget Limit", "Status"]
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Trebuchet MS'
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = EMERALD
            
    for r, item in enumerate(category_analysis):
        cat_name = item['category'].capitalize()
        spent_str = f"{symbol}{item['spent']:,.2f}"
        budget_str = f"{symbol}{item['budget']:,.2f}" if item['budget'] > 0 else "—"
        
        status = "Safe"
        if item['budget'] > 0 and item['spent'] > item['budget']:
            status = "Over Limit"
            
        vals = [cat_name, spent_str, budget_str, status]
        for c, val in enumerate(vals):
            cell = table.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Calibri'
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                if c == 3:
                    p.font.bold = True
                    p.font.color.rgb = EMERALD if status == "Safe" else RGBColor(239, 68, 68)

    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    buffer.close()
    return pptx_bytes

if __name__ == '__main__':
    app.run(debug=True, port=5000)
